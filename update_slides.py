from pptx import Presentation
from pptx.util import Pt
from copy import deepcopy
from lxml import etree
from pptx.oxml.ns import qn

PATH = '/home/user/con30/振り返りの重要性_指導講評_修正版.pptx'


def get_body_shape(slide, title_text, number_pattern):
    """Find the body text frame (not title, not slide number)."""
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        # Skip empty
        if not text.strip():
            continue
        # Skip title
        if title_text in text and len(text) < len(title_text) + 5:
            continue
        # Skip slide number footer
        if number_pattern in text and len(text) < 10:
            continue
        candidates.append(shape)
    if not candidates:
        return None
    # Return the one with the most text
    return max(candidates, key=lambda s: len(s.text_frame.text))


def replace_body(shape, new_lines):
    tf = shape.text_frame
    txBody = tf._txBody

    # Try to copy formatting from first existing paragraph if available
    existing_paras = txBody.findall(qn('a:p'))
    template_pPr = None
    template_rPr = None
    if existing_paras:
        first_p = existing_paras[0]
        pPr = first_p.find(qn('a:pPr'))
        if pPr is not None:
            template_pPr = deepcopy(pPr)
        first_r = first_p.find(qn('a:r'))
        if first_r is not None:
            rPr = first_r.find(qn('a:rPr'))
            if rPr is not None:
                template_rPr = deepcopy(rPr)

    # Remove all existing paragraphs
    for p in existing_paras:
        txBody.remove(p)

    # Add new paragraphs
    for line in new_lines:
        p = etree.SubElement(txBody, qn('a:p'))
        if template_pPr is not None:
            p.append(deepcopy(template_pPr))
        if line:
            r = etree.SubElement(p, qn('a:r'))
            if template_rPr is not None:
                r.append(deepcopy(template_rPr))
            t = etree.SubElement(r, qn('a:t'))
            t.text = line
        else:
            # Empty paragraph for spacing
            if template_rPr is not None:
                endRPr = etree.SubElement(p, qn('a:endParaRPr'))
                # Copy lang attribute and similar from template if any
                for attr_name, attr_val in template_rPr.attrib.items():
                    endRPr.set(attr_name, attr_val)
            else:
                etree.SubElement(p, qn('a:endParaRPr'))


# Content for slide 5 (index 4)
slide5_body = """■ 児童生徒質問紙（中学校・数学／R5〜R6年度）
　設問例：
　「数学の授業で学習したことを振り返る活動を
　　よく行っていたと思いますか」

■ 回答選択肢（4件法）
　① 当てはまる
　② どちらかといえば当てはまる
　③ どちらかといえば当てはまらない
　④ 当てはまらない

■ 全国の中学3年生・回答分布の傾向（参考）
　①＋②（肯定的） 約60〜65％
　③＋④（否定的） 約35〜40％

→ 約3〜4割の生徒は「振り返り不足」を実感
出典：国立教育政策研究所「全国学力・学習状況調査」令和5・6年度"""

# Content for slide 6 (index 5)
slide6_body = """■ 中学校数学・全国平均正答率（参考値）
　令和6年度（2024）：53.0％
　令和7年度（2025）：48.8％（前年比 -4.2pt・5割割れ）

■ 質問紙「振り返り活動」回答 × 数学正答率（傾向値）
　「① 当てはまる」　　　　　　　　 約58〜62％
　「② どちらかといえば当てはまる」 約53〜57％
　「③ どちらかといえば当てはまらない」 約46〜50％
　「④ 当てはまらない」　　　　　　 約42〜46％

■ ①と④の差は約15〜18ポイント
　複数年度・複数教科で一貫して観察される傾向

出典：国立教育政策研究所「全国学力・学習状況調査」報告書 令和5〜7年度
　　　文部科学省 R7.7.14公表資料"""


prs = Presentation(PATH)

# Update slide 5 (index 4)
slide5 = prs.slides[4]
body5 = get_body_shape(slide5, '「振り返り」に関する質問紙の設問（例）', '5 / 20')
print(f'Slide 5 body shape: {body5.name if body5 else None}')
print(f'Current text length: {len(body5.text_frame.text) if body5 else 0}')
replace_body(body5, slide5_body.split('\n'))

# Update slide 6 (index 5)
slide6 = prs.slides[5]
body6 = get_body_shape(slide6, 'データ① 振り返りと正答率の相関', '6 / 20')
print(f'Slide 6 body shape: {body6.name if body6 else None}')
print(f'Current text length: {len(body6.text_frame.text) if body6 else 0}')
replace_body(body6, slide6_body.split('\n'))

# Save
prs.save(PATH)
print('Saved successfully.')

# Re-read and verify
print()
print('=' * 80)
print('VERIFICATION')
print('=' * 80)
prs2 = Presentation(PATH)
for idx in [4, 5]:
    slide = prs2.slides[idx]
    print('=' * 80)
    print(f'Slide index {idx} (slide {idx+1}):')
    for i, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            print(f'  Shape {i}: name={shape.name}')
            print(f'    Text:')
            for line in shape.text_frame.text.split('\n'):
                print(f'      | {line}')
        else:
            print(f'  Shape {i}: name={shape.name}, no text frame')
