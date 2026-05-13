#!/usr/bin/env python3
"""
Reorganize PowerPoint: 42 slides → 20 slides, reorder, update 目次 and slide numbers.
"""

import shutil
import os
from copy import deepcopy
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

SOURCE = "/root/.claude/uploads/c4cd298c-6e58-4b7a-a212-f856f08e93c2/a05214a1-080520_____________.pptx"
OUTPUT = "/home/user/con30/振り返りの重要性_指導講評_修正版.pptx"

# ── helpers ──────────────────────────────────────────────────────────────────

def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    slide = prs.slides[index]
    slide_part = slide.part
    for rId, rel in list(prs.part.rels.items()):
        try:
            if rel.target_part == slide_part:
                prs.part.drop_rel(rId)
                break
        except Exception:
            pass
    sldIdLst.remove(sldIdLst[index])


def reorder_slides(prs, new_order):
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for i in new_order:
        sldIdLst.append(slides[i])


def get_slide_title(slide):
    """Return first non-empty text found on a slide (up to 40 chars)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                return txt[:40]
    return "(no text)"


def replace_in_slide(slide, replacements):
    """Replace text strings in all runs of a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)


# ── Step 1: Copy ─────────────────────────────────────────────────────────────

print("Step 1: Copying source to output …")
shutil.copy2(SOURCE, OUTPUT)
print(f"  Copied → {OUTPUT}")

# ── Step 2: Delete slides ─────────────────────────────────────────────────────

print("\nStep 2: Deleting unwanted slides …")
prs = Presentation(OUTPUT)
print(f"  Slide count before deletion: {len(prs.slides)}")

to_delete = sorted([41, 40, 39, 38, 37, 36, 35, 34, 33, 32, 31,
                    26, 21, 19, 18, 17, 16, 15, 13, 11, 10, 3],
                   reverse=True)
print(f"  Deleting indices (high→low): {to_delete}")

for idx in to_delete:
    try:
        delete_slide(prs, idx)
        print(f"    Deleted index {idx}")
    except Exception as e:
        print(f"    WARNING: could not delete index {idx}: {e}")

print(f"  Slide count after deletion: {len(prs.slides)}")

print("\n  Slide titles after deletion:")
expected_labels = [
    "表紙", "目次", "問いかけ", "子供たちの実態（数値）", "質問紙設問",
    "データ①", "データ②", "なぜ相関するか",
    "本時①評価手段", "本時②視覚化", "本時③話合い", "よかった点・提案",
    "H27→R7論点整理", "H27論点整理", "R7論点整理",
    "H29学習指導要領", "よい振り返り", "振り返り指導5ポイント",
    "まとめ", "参考文献",
]
for i, slide in enumerate(prs.slides):
    label = expected_labels[i] if i < len(expected_labels) else "?"
    print(f"    [{i:02d}] {label:20s} | {get_slide_title(slide)}")

# ── Step 3: Reorder ───────────────────────────────────────────────────────────

print("\nStep 3: Reordering slides …")
new_order = [0, 1, 2, 3, 4, 5, 6, 7, 12, 13, 14, 15, 16, 17, 8, 9, 10, 11, 18, 19]
reorder_slides(prs, new_order)
print(f"  Reorder applied: {new_order}")
print(f"  Slide count after reorder: {len(prs.slides)}")

# ── Step 4: Update 目次 slide ─────────────────────────────────────────────────

print("\nStep 4: Updating 目次 slide (index 1) …")
slide = prs.slides[1]
updated = False
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    tf = shape.text_frame
    full_text = tf.text
    if '研究主題' in full_text or '国語科' in full_text or '乙藤' in full_text:
        print(f"  Found target shape: [{full_text[:60].strip()}]")
        txBody = tf._txBody
        # Remove all existing <a:p> elements
        for p in txBody.findall(qn('a:p')):
            txBody.remove(p)

        new_lines = [
            '① 全国学力・学習状況調査から見える',
            '　「振り返り」と学力の相関関係',
            '② 文科省資料から読み解く',
            '　「振り返り」の重要性（H27・R7論点整理）',
            '③ 本時の授業について',
        ]

        for line in new_lines:
            p_xml = (
                '<a:p xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:r><a:t>{line}</a:t></a:r></a:p>'
            )
            p_elem = etree.fromstring(p_xml)
            txBody.append(p_elem)

        print("  目次 text replaced successfully.")
        updated = True
        break

if not updated:
    print("  WARNING: target shape for 目次 not found — checking all shapes on slide 1:")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(f"    Shape text: [{shape.text_frame.text[:80].strip()}]")

# ── Step 5: Update slide numbers ──────────────────────────────────────────────

print("\nStep 5: Updating slide number strings …")
replacements = {
    "3 / 21":  "3 / 20",
    "6 / 21":  "5 / 20",
    "7 / 21":  "6 / 20",
    "8 / 21":  "7 / 20",
    "9 / 21":  "8 / 20",
    "11 / 21": "9 / 20",
    "12 / 21": "10 / 20",
    "13 / 21": "11 / 20",
    "14 / 21": "12 / 20",
    "15 / 21": "13 / 20",
    "16 / 21": "13 / 20",
    "17 / 21": "14 / 20",
    "20 / 21": "18 / 20",
    "21 / 21": "19 / 20",
    "/ 21":    "/ 20",
}

for i, slide in enumerate(prs.slides):
    try:
        replace_in_slide(slide, replacements)
    except Exception as e:
        print(f"  WARNING: slide {i} replacement error: {e}")

print("  Slide number replacement done.")

# ── Step 6: Save & verify ─────────────────────────────────────────────────────

print("\nStep 6: Saving …")
prs.save(OUTPUT)
print(f"  Saved → {OUTPUT}")

size_kb = os.path.getsize(OUTPUT) / 1024
print(f"  File size: {size_kb:.1f} KB")
print(f"  Total slides: {len(prs.slides)}")

print("\n  Final slide listing:")
for i, slide in enumerate(prs.slides):
    print(f"    [{i+1:02d}] {get_slide_title(slide)}")

print("\nDone.")
