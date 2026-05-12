#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color constants
NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ORANGE = RGBColor(0xC5, 0x5A, 0x11)
LIGHT_BLUE_BG = RGBColor(0xD6, 0xE4, 0xF7)

# Font preferences
JP_FONT = "BIZ UDPGothic"

def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_name=JP_FONT, font_size=20, bold=False,
                color=DARK_GRAY, align=PP_ALIGN.LEFT,
                word_wrap=True, line_spacing=None):
    """Add a textbox with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txBox, tf

def add_multiline_textbox(slide, lines, left, top, width, height,
                           font_name=JP_FONT, font_size=18,
                           color=DARK_GRAY, align=PP_ALIGN.LEFT,
                           line_spacing_pt=None, bold=False):
    """Add a textbox with multiple paragraphs."""
    from pptx.oxml.ns import qn
    from lxml import etree

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = align

        # Handle font size override per line
        fsize = font_size
        fbold = bold
        fcolor = color

        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(fsize)
        run.font.bold = fbold
        run.font.color.rgb = fcolor

        if line_spacing_pt:
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            lnSpc = pPr.find(qn('a:lnSpc'))
            if lnSpc is None:
                lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts_elem = lnSpc.find(qn('a:spcPts'))
            if spcPts_elem is None:
                spcPts_elem = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts_elem.set('val', str(int(line_spacing_pt * 100)))

    return txBox, tf

def add_rect(slide, left, top, width, height, fill_color):
    """Add a filled rectangle."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_slide_number(slide, prs, slide_num, total=21):
    """Add slide number at bottom right."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Cm(0.5)
    w = Cm(2)
    h = Cm(0.6)
    left = slide_w - w - margin
    top = slide_h - h - Cm(0.2)
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_num} / {total}"
    run.font.name = JP_FONT
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

def add_content_title(slide, title_text, prs):
    """Add a navy title bar at top of content slide."""
    slide_w = prs.slide_width
    bar_h = Cm(1.6)
    rect = add_rect(slide, 0, 0, slide_w, bar_h, NAVY)

    # Title text in the bar
    margin = Cm(0.5)
    txBox = slide.shapes.add_textbox(margin, Cm(0.15), slide_w - margin*2, bar_h - Cm(0.3))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = JP_FONT
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return bar_h

def make_title_slide(prs, slide_layout, title, subtitle, footer):
    """Create title slide (Slide 1) with navy background."""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Main title
    add_textbox(slide, title,
                Cm(2), Cm(5),
                slide_w - Cm(4), Cm(3),
                font_size=40, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, subtitle,
                Cm(2), Cm(8.5),
                slide_w - Cm(4), Cm(2),
                font_size=24, bold=False, color=RGBColor(0xBF, 0xD7, 0xFF),
                align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, footer,
                Cm(1), slide_h - Cm(2.2),
                slide_w - Cm(2), Cm(1.5),
                font_size=13, bold=False, color=RGBColor(0xBF, 0xD7, 0xFF),
                align=PP_ALIGN.CENTER)

    return slide

def make_section_slide(prs, slide_layout, title, subtitle=None):
    """Create section header slide with blue background."""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BLUE)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if subtitle:
        # Title (e.g. "第１部")
        add_textbox(slide, title,
                    Cm(2), Cm(4),
                    slide_w - Cm(4), Cm(2.5),
                    font_size=36, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)
        # Subtitle
        add_textbox(slide, subtitle,
                    Cm(2), Cm(7),
                    slide_w - Cm(4), Cm(4),
                    font_size=26, bold=False, color=RGBColor(0xD6, 0xE4, 0xFF),
                    align=PP_ALIGN.CENTER)
    else:
        # Single title centered
        add_textbox(slide, title,
                    Cm(2), Cm(5.5),
                    slide_w - Cm(4), Cm(5),
                    font_size=32, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER)

    return slide

def make_content_slide(prs, slide_layout, title, body_lines, slide_num):
    """Create content slide with white background and navy title bar."""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    bar_h = add_content_title(slide, title, prs)

    # Body text
    body_top = bar_h + Cm(0.4)
    body_h = slide_h - body_top - Cm(1.2)
    body_left = Cm(0.8)
    body_width = slide_w - Cm(1.6)

    add_multiline_textbox(slide, body_lines,
                          body_left, body_top,
                          body_width, body_h,
                          font_size=17, color=DARK_GRAY,
                          line_spacing_pt=22)

    add_slide_number(slide, prs, slide_num)
    return slide

def make_last_title_slide(prs, slide_layout, title, body_lines, closing, footer, slide_num):
    """Create final title slide (Slide 21) with navy background."""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, NAVY)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Title
    add_textbox(slide, title,
                Cm(2), Cm(0.8),
                slide_w - Cm(4), Cm(1.8),
                font_size=34, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    # Divider line
    from pptx.util import Pt as PtU
    line_shape = slide.shapes.add_shape(1, Cm(3), Cm(2.8), slide_w - Cm(6), Cm(0.05))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0x7F, 0xA7, 0xD8)
    line_shape.line.fill.background()

    # Body lines
    add_multiline_textbox(slide, body_lines,
                          Cm(2), Cm(3.1),
                          slide_w - Cm(4), Cm(5),
                          font_size=18, color=RGBColor(0xD6, 0xE4, 0xFF),
                          line_spacing_pt=26)

    # Closing message (orange, larger, centered)
    add_textbox(slide, closing,
                Cm(1.5), Cm(9.5),
                slide_w - Cm(3), Cm(3.5),
                font_size=22, bold=True, color=ORANGE,
                align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, footer,
                Cm(1.5), slide_h - Cm(1.8),
                slide_w - Cm(3), Cm(1.2),
                font_size=13, bold=False, color=RGBColor(0x8F, 0xAF, 0xD8),
                align=PP_ALIGN.CENTER)

    add_slide_number(slide, prs, slide_num)
    return slide


def main():
    prs = Presentation()

    # Set widescreen 16:9
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    # -------------------------
    # Slide 1: Title slide
    # -------------------------
    make_title_slide(prs, blank_layout,
        "振り返りの重要性",
        "―学びを深める省察の力―",
        "令和８年５月20日（水）｜練馬区立豊玉中学校 第５校時 数学科授業参観後｜練馬区教育委員会 指導主事"
    )

    # -------------------------
    # Slide 2: Section - 本日の流れ
    # -------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, BLUE)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    add_textbox(slide2, "本日の流れ",
                Cm(2), Cm(1.2),
                slide_w - Cm(4), Cm(1.8),
                font_size=30, bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)

    flow_lines = [
        "① 全国学力・学習状況調査から見える",
        "　「振り返り」と「学力」の相関関係",
        "",
        "② 文科省資料から読み解く",
        "　「振り返り」の重要性",
        "　（H27・R7 論点整理ほか）",
        "",
        "③ 本時の授業について",
    ]
    add_multiline_textbox(slide2, flow_lines,
                          Cm(3), Cm(3.5),
                          slide_w - Cm(5), slide_h - Cm(5),
                          font_size=20, color=RGBColor(0xE8, 0xF0, 0xFF),
                          line_spacing_pt=26)
    add_slide_number(slide2, prs, 2)

    # -------------------------
    # Slide 3: はじめに
    # -------------------------
    body3 = [
        "あなたの授業の「振り返り」は何分ありますか？",
        "",
        "どんな内容を書かせていますか？",
        "",
        "それは、子供の「学び」につながっていますか？",
    ]
    make_content_slide(prs, blank_layout,
                       "はじめに―問いかけ―",
                       body3, 3)

    # -------------------------
    # Slide 4: Section header - 第１部
    # -------------------------
    make_section_slide(prs, blank_layout,
                       "第１部",
                       "全国学力・学習状況調査から見える\n「振り返り」と「学力」の相関")
    # Add slide number
    s4 = prs.slides[-1]
    add_slide_number(s4, prs, 4)

    # -------------------------
    # Slide 5: 全国学力・学習状況調査とは
    # -------------------------
    body5 = [
        "■ 実施主体：文部科学省・国立教育政策研究所",
        "■ 対象：小学校６年生・中学校３年生",
        "■ 調査内容",
        "　・教科に関する調査（国語・数学/算数・英語 等）",
        "　・学習意欲・学習方法・学習環境に関する質問紙調査",
        "■ 注目ポイント",
        "　質問紙調査 × 教科学力のクロス分析が可能",
        "　→ 授業中の活動と学力の相関が可視化できる！",
    ]
    make_content_slide(prs, blank_layout,
                       "全国学力・学習状況調査とは",
                       body5, 5)

    # -------------------------
    # Slide 6: 質問紙の設問
    # -------------------------
    body6 = [
        "質問紙における「振り返り」関連設問（中学生・数学）",
        "",
        "Q.「数学の授業で学習したことを振り返る活動を",
        "　　よく行っていた」",
        "",
        "　① 当てはまる",
        "　② どちらかといえば当てはまる",
        "　③ どちらかといえば当てはまらない",
        "　④ 当てはまらない",
        "",
        "→ この回答と数学の得点を掛け合わせると…",
        "（出典：全国学力・学習状況調査 質問紙 令和３〜５年度）",
    ]
    make_content_slide(prs, blank_layout,
                       "「振り返り」に関する質問紙の設問（例）",
                       body6, 6)

    # -------------------------
    # Slide 7: データ① 振り返りと正答率
    # -------------------------
    body7 = [
        "【グラフ挿入箇所】",
        "振り返りを「当てはまる」と答えた生徒の方が",
        "数学の正答率が明らかに高い傾向がある",
        "",
        "傾向のイメージ：",
        "　「当てはまる」      → 正答率 高",
        "　「当てはまらない」  → 正答率 低",
        "",
        "■ この差は偶然ではない",
        "　複数年度・複数教科で一貫して確認される傾向",
        "（出典：国立教育政策研究所 報告書 令和３〜５年度）",
    ]
    make_content_slide(prs, blank_layout,
                       "データ① 振り返りと正答率の相関",
                       body7, 7)

    # -------------------------
    # Slide 8: データ② 振り返りは習慣になっているか
    # -------------------------
    body8 = [
        "■ 課題：「振り返りを行った」と感じていない生徒が一定数存在",
        "",
        "■ 教師側の認識と生徒の実感の乖離",
        "　教師：「振り返りをさせている」",
        "　生徒：「振り返った実感がない」",
        "",
        "■ 問い直し：",
        "　形式的な振り返りになっていないか？",
        "　ただ「書かせる」だけになっていないか？",
        "",
        "→ 振り返りの「量」より「質」が問われている",
    ]
    make_content_slide(prs, blank_layout,
                       "データ② 振り返りは習慣になっているか",
                       body8, 8)

    # -------------------------
    # Slide 9: なぜ振り返りが学力と相関するのか
    # -------------------------
    body9 = [
        "■ メタ認知（Metacognition）の観点から",
        "",
        "　振り返り ＝ 自分の思考・理解を「対象化」する行為",
        "",
        "　　モニタリング：「わかったか？」を確認する",
        "　　コントロール：「どう学び直すか」を調整する",
        "",
        "■ J.ハッティ「可視化された学習」（2009）より",
        "　　メタ認知的方略の効果量：d = 0.60（高い効果）",
        "",
        "→ 振り返りは「学び方を学ぶ」ことにつながる",
    ]
    make_content_slide(prs, blank_layout,
                       "なぜ振り返りが学力と相関するのか",
                       body9, 9)

    # -------------------------
    # Slide 10: Section header - 第２部
    # -------------------------
    make_section_slide(prs, blank_layout,
                       "第２部",
                       "文科省資料から読み解く\n「振り返り」の重要性")
    s10 = prs.slides[-1]
    add_slide_number(s10, prs, 10)

    # -------------------------
    # Slide 11: H27→R7論点整理
    # -------------------------
    body11 = [
        "■ ２つの「論点整理」が示す一貫したメッセージ",
        "",
        "　平成27年（2015）「論点整理」",
        "　　↓ 学習指導要領改訂（H29）へ",
        "　令和３年（2021）「令和の日本型学校教育」答申",
        "　　↓ さらなる深化へ",
        "　令和７年（2025）「論点整理」",
        "　　↓ 次期学習指導要領改訂（R8〜9年告示予定）へ",
        "",
        "■ 一貫して求められていること：",
        "　「子供が自分の学びを調整できる力」",
        "　　　　振り返りはその根幹",
    ]
    make_content_slide(prs, blank_layout,
                       "H27「論点整理」→ R7「論点整理」の流れ",
                       body11, 11)

    # -------------------------
    # Slide 12: H27論点整理の核心
    # -------------------------
    body12 = [
        "■ アクティブ・ラーニングの３つの視点",
        "",
        "　① 深い学び",
        "　　→「見方・考え方」を働かせた習得・活用・探究",
        "　② 対話的な学び",
        "　　→ 他者との協働・対話を通じた思考の広がり",
        "　③ 主体的な学び　←ここに「振り返り」が明記",
        "　　→「学習を自己調整しながら学ぼうとする態度」",
        "",
        "■ キーワード：「自己調整」",
        "　→ 振り返りなしに、自己調整は生まれない",
    ]
    make_content_slide(prs, blank_layout,
                       "H27「論点整理」の核心",
                       body12, 12)

    # -------------------------
    # Slide 13: R7論点整理
    # -------------------------
    body13 = [
        "■ 自己調整学習のさらなる重視",
        "　「見通し → 学習 → 振り返り」の往還を",
        "　すべての教科・場面で意図的に位置づける",
        "",
        "■ ウェルビーイングとの接続",
        "　振り返りが「学びへの自信・有能感」を育む",
        "　→ 学ぶことが楽しい、という実感へ",
        "",
        "■ デジタルを活用した振り返りの蓄積",
        "　１人１台端末を活用し、振り返りを「見える化」",
        "　→ 学習履歴として蓄積・次の学びへ接続",
        "",
        "■ H27から10年。振り返りの重要性はより深化している。",
    ]
    make_content_slide(prs, blank_layout,
                       "R7「論点整理」が示す新たな強調点",
                       body13, 13)

    # -------------------------
    # Slide 14: 学習指導要領（H29）
    # -------------------------
    body14 = [
        "■ 総則 第３の１（主体的・対話的で深い学び）",
        "",
        "　「学習の見通しを立てたり学習したことを振り返ったり",
        "　　する活動を、計画的に取り入れ…」",
        "",
        "■ 数学科の目標（中学校学習指導要領）",
        "",
        "　「問題解決の過程を振り返って評価・改善しようとする",
        "　　態度を養う」",
        "",
        "■ 重要ポイント",
        "　→ 振り返りは「目標」に明記された重要な学習活動",
        "　→「させてあげること」ではなく「育てるべき力」",
    ]
    make_content_slide(prs, blank_layout,
                       "学習指導要領（H29）における「振り返り」",
                       body14, 14)

    # -------------------------
    # Slide 15: 国際的潮流
    # -------------------------
    body15 = [
        "■ OECD ラーニング・コンパス 2030",
        "",
        "　「エージェンシー（主体性）」育成の中核に",
        "　「振り返り（Reflection）」を位置づけ",
        "",
        "　　見通し（Anticipation）",
        "　　　　↓",
        "　　行動（Action）",
        "　　　　↓",
        "　　振り返り（Reflection）← ここが核心！",
        "",
        "■「自律した学習者」育成への国際的合意",
        "　→ 日本の授業でも意図的・系統的な指導が必要",
    ]
    make_content_slide(prs, blank_layout,
                       "国際的潮流（OECD Education 2030）",
                       body15, 15)

    # -------------------------
    # Slide 16: よい振り返りとは
    # -------------------------
    body16 = [
        "■ 振り返りの「質」を問う",
        "",
        "　✕ 低質な振り返りの例：",
        "　　「今日の授業は楽しかった」",
        "　　「二次方程式がわかった」",
        "　　「難しかった」",
        "",
        "　○ 高質な振り返りの例：",
        "　　「因数分解でとけない場合に解の公式を使う、という",
        "　　　判断ができるようになった。次回は使う場面を",
        "　　　自分で判断できるようにしたい」",
        "",
        "■ 振り返りの３つの問い（軸）",
        "　① 何がわかったか（理解）　② なぜそう考えたか（根拠）　③ 次に何をしたいか（見通し）",
    ]
    make_content_slide(prs, blank_layout,
                       "「よい振り返り」とはどんなものか",
                       body16, 16)

    # -------------------------
    # Slide 17: 振り返り指導の５つのポイント
    # -------------------------
    body17 = [
        "① 振り返りの「視点」を事前に示す",
        "　→ 本時のめあてと対応させて板書する",
        "",
        "② 「書く」ことで思考を外化させる",
        "　→ 頭の中だけでは曖昧なままになりやすい",
        "",
        "③ 振り返りを「見取る」教師の目をもつ",
        "　→ 次時の指導改善・個への対応に活用",
        "",
        "④ 振り返りを「共有」する場を設ける",
        "　→ 他者の振り返りから学ぶ機会に",
        "",
        "⑤ 振り返りを「累積・活用」する仕組みをつくる",
        "　→ 前時の振り返りを次時の導入へ接続",
    ]
    make_content_slide(prs, blank_layout,
                       "振り返り指導の５つのポイント",
                       body17, 17)

    # -------------------------
    # Slide 18: Section header - 第３部
    # -------------------------
    make_section_slide(prs, blank_layout,
                       "第３部",
                       "本時の授業について")
    s18 = prs.slides[-1]
    add_slide_number(s18, prs, 18)

    # -------------------------
    # Slide 19: 本時の授業の概要
    # -------------------------
    body19 = [
        "■ 日　時：令和８年５月20日（水）第５校時",
        "■ 学　校：練馬区立豊玉中学校",
        "■ 学年学級：○年○組（○名）",
        "■ 授業者：○○ ○○ 先生",
        "■ 単 元 名：○○○○○○○○○○",
        "■ 本時の目標：",
        "　「○○○○○○○○○○○○○○○○○○」",
        "",
        "■ 本時の学習活動の流れ",
        "　導入（○分）→ 展開（○分）→ まとめ・振り返り（○分）",
        "",
        "※ 参観当日に記入",
    ]
    make_content_slide(prs, blank_layout,
                       "本時の授業の概要",
                       body19, 19)

    # -------------------------
    # Slide 20: よかった点・提案
    # -------------------------
    body20 = [
        "【よかった点】",
        "　・めあてと振り返りが明確に対応していた",
        "　・振り返りの時間がしっかり確保されていた",
        "　・「なぜそう考えたか」まで書いている生徒が見られた",
        "",
        "【さらなる充実のための提案】",
        "　提案① 振り返りの「視点」の明示",
        "　　→ めあてに対応した振り返りの問いを板書する",
        "　提案② 振り返りの「共有」の工夫",
        "　　→ 数名の振り返りを紹介し、学び合いに活用",
        "　提案③ 振り返りの「蓄積と活用」",
        "　　→ 前時の振り返りを次時の導入につなげる",
    ]
    make_content_slide(prs, blank_layout,
                       "本時授業のよかった点・さらなる提案",
                       body20, 20)

    # -------------------------
    # Slide 21: Final title slide
    # -------------------------
    body21 = [
        "① 「振り返り」のある授業は、学力と相関している（データ）",
        "② 「振り返り」は学習指導要領・答申が求める",
        "　　「主体的な学び」の核心である（理論）",
        "③ 「振り返り」は、形式ではなく質が問われる（実践）",
    ]
    closing21 = "「一時間一時間の振り返りが、\n　子供の一生の学び方をつくる」"
    footer21 = "今日の授業から、また明日の授業改善へ。ともに学び続けましょう。"

    make_last_title_slide(prs, blank_layout,
                          "まとめ",
                          body21, closing21, footer21, 21)

    # Save
    output_path = "/home/user/con30/振り返りの重要性_指導講評.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")

    import os
    size = os.path.getsize(output_path)
    print(f"File size: {size:,} bytes ({size/1024:.1f} KB)")
    print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
